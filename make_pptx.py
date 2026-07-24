import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Colors
    NAVY = RGBColor(15, 23, 42)
    INDIGO = RGBColor(79, 70, 229)
    DARK_GRAY = RGBColor(51, 65, 85)
    LIGHT_GRAY = RGBColor(241, 245, 249)
    CARD_BG = RGBColor(248, 250, 252)
    BORDER_COLOR = RGBColor(226, 232, 240)
    GREEN = RGBColor(16, 185, 129)
    WHITE = RGBColor(255, 255, 255)

    def add_header(slide, tag_text, title_text):
        # Top Category Tag
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = tag_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = INDIGO

        # Main Title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = NAVY

    # --- SLIDE 1: Title / Cover ---
    s1 = prs.slides.add_slide(blank_layout)
    
    # Title Box
    tb = s1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "ACADEMIC PROTOTYPE — VIVA PRESENTATION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = INDIGO
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "UniStay"
    p2.font.size = Pt(48)
    p2.font.bold = True
    p2.font.color.rgb = NAVY
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Boarding Selection & Management System"
    p3.font.size = Pt(22)
    p3.font.bold = True
    p3.font.color.rgb = INDIGO
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = "A Web-Based Serverless Platform for Off-Campus Student Accommodation\nat Sabaragamuwa University of Sri Lanka (SUSL)"
    p4.font.size = Pt(14)
    p4.font.color.rgb = DARK_GRAY
    p4.alignment = PP_ALIGN.CENTER

    p5 = tf.add_paragraph()
    p5.text = "\nPresented by: Saliya Wickramasingha"
    p5.font.size = Pt(13)
    p5.font.bold = True
    p5.font.color.rgb = NAVY
    p5.alignment = PP_ALIGN.CENTER


    # --- SLIDE 2: Problem Statement ---
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "Problem Statement", "Why Does UniStay Exist?")
    
    problems = [
        ("Manual House Hunting", "Students physically walk through Pambahinna, Kumbalgama & nearby areas looking for vacant boarding signs, causing extreme stress and waste of time."),
        ("Pricing Transparency Issues", "Lack of a centralized directory means owners charge arbitrary prices with zero benchmarks or student-side visibility."),
        ("Unknown Distance to Faculty", "Students have no way to verify how far a property is from their specific academic faculty building before moving in."),
        ("No Direct Communication", "Heavy reliance on brokers and middlemen leads to exploitation, extra costs, and delayed information flow.")
    ]

    for i, (p_title, p_desc) in enumerate(problems):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.6 + row * 2.6)
        
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {p_title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = p_desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_GRAY


    # --- SLIDE 3: Objectives ---
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Project Objectives", "What UniStay Achieves")

    objectives = [
        "Develop a centralized web platform connecting students and boarding owners directly.",
        "Calculate precise geographic distances from any boarding place to 7 SUSL faculties using coordinate formulas.",
        "Provide advanced search, filtering, and sorting by price, gender rules (Boys/Girls), and proximity.",
        "Implement a secure, role-based access control system separating student and owner dashboards.",
        "Establish a direct communication channel via in-app messaging without middlemen.",
        "Enable students to bookmark, save, and compare shortlisted properties on a personalized dashboard.",
        "Provide boarding owners with an intuitive CRUD panel to manage listings, photos, and map coordinates."
    ]

    tb = s3.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, obj in enumerate(objectives):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•   {obj}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)


    # --- SLIDE 4: Architecture ---
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "System Architecture", "Serverless BaaS Architecture")

    arch_steps = [
        ("Web Browser", "HTML5, Vanilla CSS3, ES6 Modules"),
        ("Firebase SDK", "firebase-config.js & navbar.js"),
        ("Firebase Auth", "Email/Password & Google SSO"),
        ("Realtime DB", "Cloud JSON NoSQL Storage")
    ]

    for i, (a_title, a_desc) in enumerate(arch_steps):
        left = Inches(0.8 + i * 2.95)
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), Inches(2.7), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = INDIGO
        
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Step {i+1}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = INDIGO
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = a_title
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = NAVY
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = f"\n{a_desc}"
        p3.font.size = Pt(11)
        p3.font.color.rgb = DARK_GRAY
        p3.alignment = PP_ALIGN.CENTER

    # Bottom notes
    tb = s4.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Architectural Advantages:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    notes = [
        "Zero server maintenance: Static hosting on GitHub Pages with direct cloud database connectivity.",
        "Real-time data synchronization for instant messaging and search filters.",
        "Client-side geospatial logic using the Haversine formula eliminating expensive API dependencies."
    ]
    for n in notes:
        p_n = tf.add_paragraph()
        p_n.text = f"•  {n}"
        p_n.font.size = Pt(12)
        p_n.font.color.rgb = DARK_GRAY


    # --- SLIDE 5: Key Features ---
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Key Features", "What The System Can Do")

    features = [
        ("Firebase Authentication", "Email/Password and Google OAuth with strict role-based page redirects."),
        ("Listing Management", "Full CRUD operations with client-side canvas image compression."),
        ("Search & Filter", "Price range sliders, gender type chips (Boys/Girls), and proximity filters."),
        ("Distance Calculator", "Auto coordinate extraction & Haversine distance calculation to 7 faculties."),
        ("Real-Time Messaging", "Direct student-owner communication thread with live toast alerts."),
        ("Favorites System", "Bookmark and compare shortlisted properties on personal student dashboard."),
        ("Interactive Map", "Leaflet.js interactive map displaying all boarding locations around campus."),
        ("Mobile Responsive", "Custom 2-column grid layout optimized for smartphone screens.")
    ]

    for i, (f_title, f_desc) in enumerate(features):
        row = i // 4
        col = i % 4
        left = Inches(0.8 + col * 2.95)
        top = Inches(1.6 + row * 2.6)

        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.7), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = f_desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_GRAY


    # --- SLIDE 6: Database Structure ---
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Database Structure", "Firebase Realtime Database JSON Schema")

    nodes = [
        ("/users/{uid}", ["full_name", "email", "role", "created_at"]),
        ("/listings/{id}", ["title", "price", "type", "lat", "lng", "phone", "owner_id", "image_url", "is_rented"]),
        ("/favorites/{uid_listingId}", ["user_id", "listing_id", "created_at"]),
        ("/messages/{id}", ["sender_id", "receiver_id", "listing_id", "content", "is_read", "created_at"])
    ]

    for i, (node_path, fields) in enumerate(nodes):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.6 + row * 2.6)

        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = node_path
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = INDIGO

        p2 = tf.add_paragraph()
        p2.text = "Fields: " + ", ".join(fields)
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_GRAY


    # --- SLIDE 7: Haversine Formula ---
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Geospatial Calculation", "Haversine Distance Formula")

    f_card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.8))
    f_card.fill.solid()
    f_card.fill.fore_color.rgb = LIGHT_GRAY
    f_card.line.color.rgb = INDIGO

    tf = f_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "a = sin²(Δlat/2) + cos(lat₁) × cos(lat₂) × sin²(Δlng/2)\nc = 2 × atan2(√a, √(1−a))\ndistance = R × c × 1.45"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    steps = [
        ("Step 1: Extract Coordinates", "Extract lat/lng from any Google Maps URL or short link using regex and unshortener proxy."),
        ("Step 2: Calculate Proximity", "Apply Haversine formula in JavaScript for all 7 university faculties."),
        ("Step 3: Render Distances", "Display exact walking distance (km) on property cards and details pages.")
    ]

    for i, (s_title, s_desc) in enumerate(steps):
        left = Inches(0.8 + i * 3.95)
        card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.7), Inches(3.7), Inches(2.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = s_desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_GRAY


    # --- SLIDE 8: System Modules & Pages ---
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "System Modules", "Application Structure by User Role")

    roles = [
        ("Student Role & Pages", [
            "index.html — Landing page & platform overview",
            "search.html — Search, filter & sort listings",
            "listing_detail.html — Detailed view & distance breakdown",
            "messages.html — Direct chat thread with landlord",
            "student_dashboard.html — Saved & bookmarked places"
        ]),
        ("Property Owner Role & Pages", [
            "dashboard.html — Manage active listings & add new place",
            "edit_listing.html — Edit property details & coordinates",
            "messages.html — View & reply to incoming student inquiries",
            "explore_map.html — Interactive map overview",
            "auth.html — Dedicated sign-in & role selection modal"
        ])
    ]

    for i, (r_title, r_pages) in enumerate(roles):
        left = Inches(0.8 + i * 5.9)
        card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(5.6), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = INDIGO if i==0 else GREEN

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = r_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY

        for pg in r_pages:
            p_pg = tf.add_paragraph()
            p_pg.text = f"•  {pg}"
            p_pg.font.size = Pt(12)
            p_pg.font.color.rgb = DARK_GRAY
            p_pg.space_after = Pt(8)


    # --- SLIDE 9: Results & Evaluation ---
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Results & Evaluation", "Expected vs. Observed Behavior")

    results = [
        ("Role-Based Access Control", "Restrict student/owner pages by role", "Works as expected"),
        ("Geospatial Distance Calculation", "Compute precise faculty proximity", "Accurate Haversine calculation"),
        ("Owner Listing Management", "Full CRUD control over properties", "Works as expected"),
        ("Real-Time Messaging Module", "Direct private student-owner chat", "Real-time delivery & notifications"),
        ("Favorites & Bookmarks", "Bookmark & compare listings on dashboard", "Works as expected"),
        ("Mobile Responsive Interface", "2-column grid layout on mobile screens", "Fully responsive")
    ]

    table_shape = s9.shapes.add_table(7, 3, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.5))
    table = table_shape.table

    headers = ["Feature", "Expected Behavior", "Observed Result"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = INDIGO
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE

    for row_idx, (f, e, o) in enumerate(results, start=1):
        table.cell(row_idx, 0).text = f
        table.cell(row_idx, 1).text = e
        table.cell(row_idx, 2).text = o
        for c in range(3):
            cell = table.cell(row_idx, c)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = NAVY

    # Summary box
    tb = s9.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Functional Correctness Rate: ~95% based on manual test execution."
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER


    # --- SLIDE 10: Proposed Enhancements ---
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "Future Work", "Proposed Enhancements")

    enhancements = [
        ("Mobile Application", "Develop a hybrid mobile app using Flutter or React Native to support background push notifications."),
        ("Payment Gateway Integration", "Integrate local payment gateways (PayHere, Stripe) allowing online room booking deposits."),
        ("Machine Learning Price Predictor", "Build a predictive data science model to estimate fair rent prices based on coordinates & amenities."),
        ("Security & Owner Verification", "Add owner identity checks and email verification to prevent fraudulent listings.")
    ]

    for i, (e_title, e_desc) in enumerate(enhancements):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.6 + row * 2.6)

        card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {e_title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = e_desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_GRAY


    # --- SLIDE 11: Conclusion ---
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "Conclusion", "Project Summary & Impact")

    conclusions = [
        "Successfully developed a functional serverless web prototype with zero backend server overhead.",
        "Implemented coordinate-based spatial analysis using the Haversine formula for 7 SUSL faculties.",
        "Delivered secure role-based access control separating student and boarding owner workflows.",
        "Enabled direct real-time messaging and live toast notification system between users.",
        "Demonstrated practical feasibility as a student welfare platform prototype for SUSL Belihuloya.",
        "Established a scalable foundation for future enhancements like mobile apps, payments, and ML pricing."
    ]

    tb = s11.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for c in conclusions:
        p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
        p.text = f"•   {c}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(14)


    # --- SLIDE 12: Q&A / Thank You ---
    s12 = prs.slides.add_slide(blank_layout)
    
    tb = s12.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = INDIGO
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "UniStay — Boarding Selection & Management System\nSabaragamuwa University of Sri Lanka"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = NAVY
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "\nQuestions & Answers"
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = DARK_GRAY
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = "Feel free to ask any questions about the system architecture, distance calculations, or implementation details."
    p4.font.size = Pt(12)
    p4.font.color.rgb = DARK_GRAY
    p4.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_path = "d:\\MyGIthub project\\unistay\\UniStay_Viva_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    build_presentation()
